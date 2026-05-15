from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


OUTPUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUTPUT_DIR / "Green_Campus_Demo_Deck.pptx"
PDF_PATH = OUTPUT_DIR / "Green_Campus_User_Guide.pdf"
MARKDOWN_PATH = OUTPUT_DIR / "DEMO_TALK_TRACK.md"

TITLE = "Green Campus Sustainability Intelligence Platform"
SUBTITLE = "User guide and explanation document"

THEME = {
    "green": RGBColor(27, 127, 98),
    "dark": RGBColor(23, 52, 45),
    "muted": RGBColor(96, 117, 111),
    "blue": RGBColor(2, 132, 199),
    "bg": RGBColor(244, 251, 247),
}


SLIDES = [
    {
        "title": "Project Overview",
        "bullets": [
            "Green Campus helps monitor energy, water, waste, and solar performance in one dashboard.",
            "It is designed for college or university campuses where usage changes with occupancy, seasons, and academic schedules.",
            "The system supports daily operations, future forecasting, sustainability scoring, and dataset import/export.",
        ],
    },
    {
        "title": "Problem We Are Solving",
        "bullets": [
            "Campus utility information is often spread across Excel files and separate teams.",
            "It is hard to understand current consumption, compare buildings, and plan for upcoming months.",
            "This platform brings all resource data into one place and converts it into decisions.",
        ],
    },
    {
        "title": "How The System Works",
        "bullets": [
            "Frontend: React dashboard for login, filters, charts, AI panels, simulator, import, reset, and export.",
            "Backend: FastAPI service that reads campus data, calculates metrics, and returns APIs for the frontend.",
            "Database: Stores raw energy, water, waste, solar, and user records.",
            "Deployment: Frontend and backend are both live, while local development is still available for new features.",
        ],
    },
    {
        "title": "Main Website Flow",
        "bullets": [
            "Step 1: User logs in to the platform.",
            "Step 2: Dashboard loads campus summary, alerts, risk, forecasts, charts, and assumptions.",
            "Step 3: User can filter by building and forecast mode.",
            "Step 4: User can import a fresh Excel file, export the current workbook, or reset data.",
            "Step 5: User reviews AI insights and runs what-if simulations.",
        ],
    },
    {
        "title": "What Data The Platform Stores",
        "bullets": [
            "Energy sheet: date, building, energy_kwh",
            "Water sheet: date, building, water_kl",
            "Waste sheet: date, building, waste_kg",
            "Solar sheet: date, building, solar_kwh",
            "User table: username, password hash, role",
        ],
    },
    {
        "title": "Dashboard Sections",
        "bullets": [
            "Operations snapshot: current scope, forecast mode, import/export controls, KPI cards, seasonal outlook",
            "Performance intelligence: carbon footprint, solar operations, resource mix, trends, alerts",
            "AI planning studio: insights, risk, forecasts, scenario simulator",
            "Governance and reference: building ranking and assumptions glossary",
        ],
    },
    {
        "title": "Calculation Logic: Summary Metrics",
        "bullets": [
            "Gross Energy = sum of electricity consumption in the selected scope",
            "Solar = sum of solar generation in the selected scope",
            "Net Energy = max(Gross Energy - Solar, 0)",
            "Export to Grid = max(Solar - Gross Energy, 0)",
            "Water = sum of water records",
            "Waste = sum of waste records",
        ],
    },
    {
        "title": "Calculation Logic: Carbon Metrics",
        "bullets": [
            "Emission factor used in the project = 0.82 kg CO2 per kWh",
            "Net Carbon = Net Energy x 0.82",
            "Gross Carbon = Gross Energy x 0.82",
            "Solar Avoided Carbon = min(Solar, Gross Energy) x 0.82",
            "This makes the carbon section easy to explain in both operational and sustainability terms.",
        ],
    },
    {
        "title": "Calculation Logic: Green Index",
        "bullets": [
            "Benchmarks: Energy 100000, Water 20000, Waste 30000",
            "Energy Score = min(Net Energy / 100000, 1)",
            "Water Score = min(Water / 20000, 1)",
            "Waste Score = min(Waste / 30000, 1)",
            "Weighted Efficiency = 0.5 x Energy Score + 0.3 x Water Score + 0.2 x Waste Score",
            "Green Index = (1 - Weighted Efficiency) x 100",
            "Higher Green Index means better sustainability performance.",
        ],
    },
    {
        "title": "How Forecasting Works",
        "bullets": [
            "Forecast modes available: daily, monthly, yearly, and seasonal.",
            "Historical records are grouped by the selected granularity.",
            "The system finds recurring cycle averages from past periods.",
            "Future values are created by repeating those cycle averages in the next periods.",
            "This gives practical operational forecasting even without heavy ML infrastructure.",
        ],
    },
    {
        "title": "Seasonal and Occupancy Logic",
        "bullets": [
            "The project includes academic occupancy factors by month.",
            "Lower occupancy months like June and July reduce expected campus demand.",
            "The seasonal panel compares current month occupancy with the next month.",
            "It also estimates when solar may offset most of the load or become export-ready.",
        ],
    },
    {
        "title": "AI Features In Plain Language",
        "bullets": [
            "Risk Engine: estimates whether upcoming usage is low, medium, or high risk.",
            "Insights: identifies anomalies, opportunities, and recommendations.",
            "Alert Center: highlights unusual changes and export-ready days.",
            "Scenario Simulator: tests what happens if energy, water, and waste reduce or solar increases.",
        ],
    },
    {
        "title": "Excel Import and Export Workflow",
        "bullets": [
            "Import Excel: uploads a fresh workbook and replaces old campus resource data.",
            "Export Excel: downloads a workbook with summary, building performance, and all raw sheets.",
            "Reset Data: clears current solar, energy, water, and waste records for a fresh view.",
            "This makes the platform practical for demos, audits, and management sharing.",
        ],
    },
    {
        "title": "Live Demo Script",
        "bullets": [
            "1. Open login page and log in",
            "2. Show dashboard overview and explain KPI cards",
            "3. Change building and forecast granularity",
            "4. Open carbon, solar, alerts, and resource charts",
            "5. Explain AI insights, seasonal outlook, and scenario simulator",
            "6. Import the sample workbook and show how data refreshes",
            "7. Export the workbook and explain the generated sheets",
        ],
    },
    {
        "title": "Why This Project Is Useful",
        "bullets": [
            "Easy for campus administrators to understand current performance.",
            "Helps compare buildings and detect wasteful patterns early.",
            "Supports planning for low-occupancy periods, solar contribution, and carbon reduction.",
            "Provides a strong base for future features like reports, alerts, optimization, and occupancy-aware forecasting.",
        ],
    },
]


GUIDE_SECTIONS = [
    (
        "1. What this system is",
        [
            "Green Campus is a dashboard that helps a campus track electricity, water, waste, solar, carbon, and future planning in one place.",
            "It is meant for a user who wants to understand the current situation clearly and make better decisions for the next period.",
        ],
    ),
    (
        "2. What the main values mean",
        [
            "Gross Energy means all electricity used before solar support is removed.",
            "Solar means electricity produced by the solar system.",
            "Net Energy means electricity still needed from the grid after using solar.",
            "Water means total water use in the selected scope.",
            "Waste means total waste in the selected scope.",
            "Carbon means estimated emissions linked to the remaining grid electricity.",
            "Green Index means the final overall sustainability score. Higher is better.",
        ],
    ),
    (
        "3. Basic formulas used by the dashboard",
        [
            "Net Energy = max(Gross Energy - Solar, 0).",
            "Export to Grid = max(Solar - Gross Energy, 0).",
            "Carbon = Net Energy x carbon factor.",
            "Gross Carbon = Gross Energy x carbon factor.",
            "Solar Avoided Carbon = min(Solar, Gross Energy) x carbon factor.",
        ],
    ),
    (
        "4. Green Index in simple words",
        [
            "The system first checks whether energy, water, and waste are high or low compared with their benchmark values.",
            "Then it combines those three results using weights.",
            "Energy has the biggest effect, water the next, and waste the least.",
            "After that, the system converts the result into a score out of 100.",
            "If energy, water, and waste are lower and more efficient, the Green Index becomes higher.",
        ],
    ),
    (
        "5. Green Index step by step",
        [
            "Step 1: Compare Net Energy with the energy benchmark.",
            "Step 2: Compare Water with the water benchmark.",
            "Step 3: Compare Waste with the waste benchmark.",
            "Step 4: Multiply those three results by their weights.",
            "Step 5: Add them to get one combined efficiency pressure value.",
            "Step 6: Convert that pressure into the final Green Index score.",
            "Simple formula: Green Index = (1 - combined efficiency pressure) x 100.",
        ],
    ),
    (
        "6. Important terms in simple language",
        [
            "Benchmark means a reference value used to judge whether usage is low or high.",
            "Weighted Efficiency means the combined effect of energy, water, and waste after giving each one different importance.",
            "Forecast Range means the system shows a likely lower and upper band instead of pretending one exact future number is guaranteed.",
            "Occupancy Effect means future demand changes because campus activity changes from month to month.",
            "Export Potential means extra solar energy that may remain after campus demand is met.",
        ],
    ),
    (
        "7. What occupancy values mean",
        [
            "1.0 means normal or full campus activity.",
            "0.9 means slightly reduced activity.",
            "0.5 means partial activity or vacation-like activity.",
            "0.2 means very low campus activity.",
            "These are activity multipliers, not exact percentages.",
        ],
    ),
    (
        "8. Forecast and seasonal logic in simple words",
        [
            "The system looks at past data and groups it by daily, monthly, yearly, or seasonal view.",
            "Then it uses those patterns to estimate future values.",
            "Lower occupancy months like June and July usually reduce expected energy and water demand.",
            "At the same time, solar may cover a larger share of demand during low-activity months.",
        ],
    ),
    (
        "9. Alerts, recommendations, and action board",
        [
            "Alerts help the user notice spikes, unusual changes, and possible problem areas.",
            "Recommendations suggest what action could improve energy, cost, carbon, or resource behavior.",
            "The Action Board lets the user mark each recommendation as suggested, planned, in progress, or completed.",
        ],
    ),
    (
        "10. Import, export, and reset",
        [
            "Import Excel loads a fresh workbook into the system.",
            "Export Excel downloads the current system data as a workbook.",
            "Reset Data clears the current records so a new dataset can be loaded.",
        ],
    ),
]


def add_bullets(text_frame, bullets):
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = THEME["dark"]
        paragraph.space_after = Pt(10)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = THEME["bg"]

    title_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(1.4))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = TITLE
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = THEME["dark"]

    subtitle_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(2.35), Inches(10.8), Inches(0.8))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = "Simple explanation of the system, website flow, calculations, and demo narrative"
    subtitle_p.font.size = Pt(18)
    subtitle_p.font.color.rgb = THEME["muted"]

    footer_box = title_slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8), Inches(0.6))
    footer_p = footer_box.text_frame.paragraphs[0]
    footer_p.text = SUBTITLE
    footer_p.font.size = Pt(16)
    footer_p.font.color.rgb = THEME["green"]

    for slide_data in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

        top_band = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(0.4)
        )
        top_band.fill.solid()
        top_band.fill.fore_color.rgb = THEME["green"]
        top_band.line.color.rgb = THEME["green"]

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.65), Inches(12), Inches(0.7))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data["title"]
        title_p.font.size = Pt(26)
        title_p.font.bold = True
        title_p.font.color.rgb = THEME["dark"]

        body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(11.7), Inches(5.35))
        add_bullets(body_box.text_frame, slide_data["bullets"])

    prs.save(PPTX_PATH)


def build_pdf():
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "TitleStyle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=22,
        textColor=colors.HexColor("#17342d"),
        spaceAfter=16,
    )
    heading_style = ParagraphStyle(
        "HeadingStyle",
        parent=styles["Heading2"],
        fontName="Helvetica-Bold",
        fontSize=15,
        textColor=colors.HexColor("#1b7f62"),
        spaceAfter=8,
        spaceBefore=8,
    )
    body_style = ParagraphStyle(
        "BodyStyle",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10.5,
        leading=15,
        textColor=colors.HexColor("#263c37"),
        spaceAfter=5,
    )

    doc = SimpleDocTemplate(
        str(PDF_PATH),
        pagesize=A4,
        rightMargin=0.55 * inch,
        leftMargin=0.55 * inch,
        topMargin=0.6 * inch,
        bottomMargin=0.6 * inch,
    )

    story = []
    story.append(Paragraph(TITLE, title_style))
    story.append(Paragraph("Simple user guide for terms, formulas, and key dashboard values", body_style))
    story.append(Spacer(1, 12))

    assumptions_table = Table(
        [
            ["Parameter", "Value"],
            ["Emission factor", "0.82 kg CO2 per kWh"],
            ["Energy cost", "8.0 per kWh"],
            ["Energy benchmark", "100000"],
            ["Water benchmark", "20000"],
            ["Waste benchmark", "30000"],
            ["Green Index weights", "Energy 0.5, Water 0.3, Waste 0.2"],
            ["Alert thresholds", "15% standard, 30% high"],
        ],
        colWidths=[2.2 * inch, 3.9 * inch],
    )
    assumptions_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1b7f62")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#c8d7d2")),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#f6faf8")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("FONTSIZE", (0, 0), (-1, -1), 9.5),
                ("LEADING", (0, 0), (-1, -1), 12),
            ]
        )
    )
    story.append(assumptions_table)
    story.append(Spacer(1, 14))

    for heading, bullets in GUIDE_SECTIONS:
        story.append(Paragraph(heading, heading_style))
        for bullet in bullets:
            story.append(Paragraph(f"• {bullet}", body_style))
        story.append(Spacer(1, 6))

    doc.build(story)


def build_markdown():
    lines = [f"# {TITLE}", "", "## Demo talking track", ""]
    for slide in SLIDES:
        lines.append(f"## {slide['title']}")
        lines.extend([f"- {bullet}" for bullet in slide["bullets"]])
        lines.append("")

    MARKDOWN_PATH.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    build_presentation()
    build_pdf()
    build_markdown()
    print(f"Created: {PPTX_PATH}")
    print(f"Created: {PDF_PATH}")
    print(f"Created: {MARKDOWN_PATH}")
